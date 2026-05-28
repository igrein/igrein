import io
import re
import logging
from PIL import Image
import pytesseract
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import pdfplumber
import cv2
import numpy as np

logger = logging.getLogger(__name__)

# ===== КОНСТАНТЫ =====
MAX_PDF_PAGES_OCR = 20
MAX_FILE_SIZE_MB = 25
OCR_TIMEOUT = 15
OCR_CONFIG = r'--oem 3 --psm 6'
OCR_FALLBACK_LANGS = ['rus+eng', 'rus', 'eng']

try:
    RESAMPLE = Image.Resampling.LANCZOS
except AttributeError:
    RESAMPLE = Image.LANCZOS


def sanitize_text(text: str) -> str:
    """ Удаляет только системный мусор.
    НЕ трогает содержимое: формулы, даты, имена, термины.
    """
    if not text:
        return ""
    
    text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', text)
    
    return text

# PARSER 
def parse_file(content: bytes, filename: str) -> str:
    """
    Главная функция парсинга.
    Универсальна для любых предметов.
    """
    if len(content) > MAX_FILE_SIZE_MB * 1024 * 1024:
        raise ValueError(f"Файл превышает {MAX_FILE_SIZE_MB} MB")
    
    filename_lower = filename.lower()
    
    # ЛОГ 1: Что загружаем
    logger.info("=" * 80)
    logger.info(f"📁 НАЧАЛО ПАРСИНГА: {filename}")
    logger.info(f"📊 РАЗМЕР ФАЙЛА: {len(content)} байт")
    logger.info("=" * 80)
    
    try:
        # 1. Извлечение текста в зависимости от формата
        if filename_lower.endswith('.txt'):
            text = _parse_txt(content)
            logger.info("📄 ФОРМАТ: TXT")
        elif filename_lower.endswith('.pdf'):
            text = _parse_pdf(content)
            logger.info("📄 ФОРМАТ: PDF")
        elif filename_lower.endswith('.docx'):
            text = _parse_docx(content)
            logger.info("📄 ФОРМАТ: DOCX")
        elif filename_lower.endswith('.pptx'):
            text = _parse_pptx(content)
            logger.info("📄 ФОРМАТ: PPTX")
        elif filename_lower.endswith(('.png', '.jpg', '.jpeg')):
            text = _parse_image(content)
            logger.info("📄 ФОРМАТ: IMAGE (OCR)")
        else:
            raise ValueError(f"Неподдерживаемый формат: {filename}")
        
        if not text or not text.strip():
            raise ValueError("Не удалось извлечь текст из файла")
        
        # ЛОГ 2: Текст ДО очистки
        logger.info("=" * 80)
        logger.info("📝 ТЕКСТ ДО SANITIZE (первые 1000 символов):")
        logger.info("=" * 80)
        logger.info(text[:1000])
        logger.info("=" * 80)
        logger.info(f"📊 ДЛИНА ТЕКСТА ДО: {len(text)} символов")
        
        # 2. Только санитайз (без агрессивных замен)
        text = sanitize_text(text)
        
        # ЛОГ 3: Текст ПОСЛЕ очистки
        logger.info("=" * 80)
        logger.info("✨ ТЕКСТ ПОСЛЕ SANITIZE (первые 1000 символов):")
        logger.info("=" * 80)
        logger.info(text[:1000])
        logger.info("=" * 80)
        logger.info(f"📊 ДЛИНА ТЕКСТА ПОСЛЕ: {len(text)} символов")
        logger.info("=" * 80)
        
        return text
        
    except Exception as e:
        logger.error(f"❌ Ошибка парсинга {filename}: {e}")
        raise


def _parse_pdf(content: bytes) -> str:
    """Парсинг PDF: сначала текст, потом OCR"""
    
    # Попытка 1: извлечь текст
    text = _extract_pdf_text(content)
    if text and len(text.strip()) > 100:
        logger.info("✅ Извлечен текст из PDF")
        return text
    
    # Попытка 2: OCR
    logger.info("📸 PDF похож на скан, используем OCR")
    text = _ocr_pdf(content)
    
    if not text or not text.strip():
        raise ValueError("Не удалось распознать PDF")
    
    return text


def _extract_pdf_text(content: bytes) -> str:
    """Извлечение текста из PDF"""
    text_parts = []
    
    # Метод 1: pdfplumber
    try:
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
        
        if text_parts:
            text = "\n".join(text_parts)
            if len(text.strip()) > 100:
                return text
    except Exception as e:
        logger.debug(f"pdfplumber: {e}")
    
    # Метод 2: PyPDF2
    text_parts = []
    try:
        reader = PdfReader(io.BytesIO(content))
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
        
        if text_parts:
            text = "\n".join(text_parts)
            if len(text.strip()) > 100:
                return text
    except Exception as e:
        logger.debug(f"PyPDF2: {e}")
    
    return ""


def _ocr_pdf(content: bytes) -> str:
    """OCR для PDF-сканов"""
    try:
        from pdf2image import convert_from_bytes
        
        # Определяем реальное количество страниц
        reader = PdfReader(io.BytesIO(content))
        total_pages = len(reader.pages)
        pages_to_ocr = min(total_pages, MAX_PDF_PAGES_OCR)
        
        logger.info(f"PDF имеет {total_pages} страниц, OCR первых {pages_to_ocr}")
        
        text_parts = []
        
        for page_num in range(1, pages_to_ocr + 1):
            try:
                images = convert_from_bytes(
                    content,
                    dpi=150,
                    first_page=page_num,
                    last_page=page_num
                )
                
                if not images:
                    continue
                
                img = images[0]
                img_processed = _preprocess_image(img)
                
                ocr_text = ""
                for lang in OCR_FALLBACK_LANGS:
                    try:
                        ocr_text = pytesseract.image_to_string(
                            img_processed,
                            lang=lang,
                            config=OCR_CONFIG,
                            timeout=OCR_TIMEOUT
                        )
                        if ocr_text and ocr_text.strip():
                            break
                    except Exception:
                        continue
                
                if ocr_text and ocr_text.strip():
                    text_parts.append(ocr_text)
                
                del img
                del img_processed
                
            except Exception as e:
                logger.warning(f"Ошибка OCR страницы {page_num}: {e}")
                continue
        
        if not text_parts:
            return ""
        
        text = "\n".join(text_parts)
        logger.info(f"✅ OCR распознал {len(text_parts)} страниц")
        return text
            
    except ImportError:
        raise ValueError("Для OCR PDF требуется установка pdf2image")
    except Exception as e:
        logger.error(f"OCR ошибка: {e}")
        return ""


def _preprocess_image(image: Image) -> Image:
    """Минимальный препроцессинг для OCR"""
    try:
        img = np.array(image)
        
        # Конвертация в оттенки серого
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY) if len(img.shape) == 3 else img
        
        # Увеличение только если слишком маленькое
        height = gray.shape[0]
        if height < 800:
            gray = cv2.resize(gray, None, fx=1.3, fy=1.3, interpolation=cv2.INTER_CUBIC)
        
        # Бинаризация
        _, binary = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
        
        return Image.fromarray(binary)
        
    except Exception as e:
        logger.warning(f"Препроцессинг: {e}")
        return image


def _parse_txt(content: bytes) -> str:
    """Парсинг TXT с определением кодировки"""
    encodings = ['utf-8', 'cp1251', 'koi8-r']
    
    for encoding in encodings:
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    
    return content.decode('utf-8', errors='ignore')

def _parse_docx(content: bytes) -> str:
    """Парсинг DOCX"""
    try:
        doc = Document(io.BytesIO(content))
        text_parts = []
        
        logger.info("📖 ПАРСИНГ DOCX:")
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                logger.debug(f"  Абзац: {text[:100]}")
                text_parts.append(text)
        
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    line = " | ".join(row_text)
                    logger.debug(f"  Таблица: {line[:100]}")
                    text_parts.append(line)
        
        if not text_parts:
            raise ValueError("DOCX не содержит текста")
        
        full_text = "\n".join(text_parts)
        logger.info(f"✅ DOCX распарсен: {len(full_text)} символов")
        
        return full_text
    
    except Exception as e:
        raise ValueError(f"Ошибка чтения DOCX: {str(e)}")


def _parse_pptx(content: bytes) -> str:
    """Парсинг PPTX"""
    try:
        prs = Presentation(io.BytesIO(content))
        text_parts = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        text_parts.append(text)
        
        if not text_parts:
            logger.warning("PPTX не содержит текста")
            return ""
        
        return "\n".join(text_parts)
    
    except Exception as e:
        logger.warning(f"Ошибка чтения PPTX: {e}")
        return ""


def _parse_image(content: bytes) -> str:
    """OCR для изображений"""
    try:
        image = Image.open(io.BytesIO(content))
        image.load()
        
        # Уменьшаем если слишком большое
        max_size = 2000
        if image.width > max_size or image.height > max_size:
            image.thumbnail((max_size, max_size), RESAMPLE)
        
        image_processed = _preprocess_image(image)
        
        text = ""
        for lang in OCR_FALLBACK_LANGS:
            try:
                text = pytesseract.image_to_string(
                    image_processed,
                    lang=lang,
                    config=OCR_CONFIG,
                    timeout=OCR_TIMEOUT
                )
                if text and text.strip():
                    break
            except Exception:
                continue
        
        if not text or not text.strip():
            raise ValueError("Не удалось распознать текст на изображении")
        
        return text
    
    except Exception as e:
        raise ValueError(f"Ошибка OCR: {str(e)}")
